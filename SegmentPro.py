"""
SegmentPro — Customer Segmentation Desktop Tool
================================================
A Windows desktop app for running psychographic / behavioral segmentation
on any survey-style Excel or CSV file.

Features
--------
• Upload any .xlsx / .xls / .csv file
• Auto-detects likely segmentation variables (numeric Likert-style items)
• Manual override to pick/exclude variables
• Cleans data (fills missing, drops empty rows, handles outliers)
• Evaluates optimal number of segments via silhouette score
• Runs K-Means clustering with reproducible results
• Auto-labels segments based on distinguishing attitudes
• Exports Excel workbook (7 sheets) and PowerPoint deck

Run:   python SegmentPro.py
Build: pyinstaller --onefile --windowed --name SegmentPro SegmentPro.py
"""

import os
import sys
import threading
import traceback
from pathlib import Path
from datetime import datetime

import numpy as np
import pandas as pd

# Tkinter imports are guarded so the engine/report classes can be imported
# in headless environments (e.g. CI). main() checks _TK_AVAILABLE.
try:
    import tkinter as tk
    from tkinter import ttk, filedialog, messagebox, scrolledtext
    _TK_AVAILABLE = True
except ImportError:
    _TK_AVAILABLE = False
    tk = ttk = filedialog = messagebox = scrolledtext = None

# ============================================================================
# CONFIGURATION
# ============================================================================
APP_NAME = "SegmentPro"
APP_VERSION = "1.0"
APP_SUBTITLE = "Customer Segmentation for Market Research"

# Visual theme
COLOR_BG = "#F4F6FA"
COLOR_PANEL = "#FFFFFF"
COLOR_PRIMARY = "#1F4E79"
COLOR_PRIMARY_DARK = "#163B5C"
COLOR_ACCENT = "#2E75B6"
COLOR_TEXT = "#1A1A1A"
COLOR_TEXT_MUTED = "#5A6773"
COLOR_BORDER = "#D6DCE5"
COLOR_SUCCESS = "#0F7B3C"
COLOR_WARN = "#B8860B"
COLOR_ERROR = "#B22222"

FONT_TITLE = ("Segoe UI", 18, "bold")
FONT_SUBTITLE = ("Segoe UI", 10)
FONT_H2 = ("Segoe UI", 12, "bold")
FONT_BODY = ("Segoe UI", 10)
FONT_SMALL = ("Segoe UI", 9)
FONT_MONO = ("Consolas", 9)

# Clustering config
MIN_K = 3
MAX_K = 7
RANDOM_STATE = 42


# ============================================================================
# UTILITIES
# ============================================================================
def smart_short(name, max_len=90):
    """Trim long survey-column names intelligently. For composite names
    with ' :: ' separators (e.g. 'Some question :: Category :: TV ads'),
    keep the tail segments where the differentiating info lives.

    Used by both the PPTX generator and the on-screen results cards so the
    user sees 'TV ads' not 'How important is each of the following as...'.
    """
    s = str(name)
    if len(s) <= max_len:
        return s
    parts = None
    for sep in (" :: ", "::"):
        if sep in s:
            parts = s.split(sep)
            break
    if parts and len(parts) > 1:
        tail = parts[-1].strip()
        for p in reversed(parts[:-1]):
            candidate = p.strip() + " · " + tail
            if len(candidate) <= max_len:
                tail = candidate
            else:
                break
        if len(tail) <= max_len:
            return "… " + tail
    if len(s) > max_len:
        return "…" + s[-(max_len - 1):]
    return s


# ============================================================================
# SEGMENTATION ENGINE (non-UI logic)
# ============================================================================
class SegmentationEngine:
    """Pure-logic segmentation pipeline. No UI dependencies."""

    def __init__(self, log_callback=None):
        self.log = log_callback or (lambda msg, level="info": None)
        self.df = None
        self.original_df = None
        self.seg_cols = []
        self.df_clean = None
        self.labels = None
        self.seg_profiles = None
        self.seg_std_profile = None
        self.seg_names = {}
        self.seg_descriptions = {}
        self.k = None
        self.silhouette_scores = {}
        self.demographics = {}

    # ---------------------------------------------------------------
    def load_file(self, filepath):
        """Load an Excel or CSV file. For Excel, use the first sheet that
        has a tabular shape (most rows)."""
        self.log(f"Loading file: {filepath}", "info")
        ext = Path(filepath).suffix.lower()

        if ext in (".xlsx", ".xls", ".xlsm"):
            xl = pd.ExcelFile(filepath)
            # Prefer a sheet named 'Data' (common convention) otherwise
            # the sheet with the most rows.
            if "Data" in xl.sheet_names:
                sheet = "Data"
            else:
                # Check each sheet's shape
                shapes = {s: pd.read_excel(filepath, sheet_name=s, nrows=0).shape[1]
                          for s in xl.sheet_names}
                row_counts = {}
                for s in xl.sheet_names:
                    try:
                        row_counts[s] = len(pd.read_excel(filepath, sheet_name=s, usecols=[0]))
                    except Exception:
                        row_counts[s] = 0
                sheet = max(row_counts, key=row_counts.get)
            self.log(f"Reading sheet: '{sheet}'", "info")
            df = pd.read_excel(filepath, sheet_name=sheet)
        elif ext == ".csv":
            # Try common encodings
            for enc in ("utf-8", "cp1252", "latin-1"):
                try:
                    df = pd.read_csv(filepath, encoding=enc, low_memory=False)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                raise ValueError("Could not decode CSV with common encodings.")
        else:
            raise ValueError(f"Unsupported file extension: {ext}")

        # Drop fully-empty rows and columns
        df = df.dropna(axis=0, how="all").dropna(axis=1, how="all").reset_index(drop=True)
        self.df = df
        self.original_df = df.copy()
        self.log(f"Loaded {len(df):,} rows × {len(df.columns):,} columns", "success")
        return df

    # ---------------------------------------------------------------
    def auto_detect_seg_cols(self):
        """Identify columns suitable for segmentation basis.

        Criteria:
        • Numeric dtype (or convertible to numeric)
        • Reasonable value range (3–10 unique values — Likert-style)
        • At least 85% non-null
        • Not an ID/count/year/date column
        """
        if self.df is None:
            return []

        df = self.df
        n = len(df)
        skip_names = {"id", "userid", "user_id", "respondent", "respondent_id",
                      "index", "date", "time", "timestamp", "year", "month",
                      "day", "elapsedtime", "username", "password"}

        candidates = []
        for col in df.columns:
            col_lower = str(col).strip().lower()
            # Skip obvious meta columns
            if any(s in col_lower for s in ("(please specify)", "other (please",
                                             "please specify")):
                continue
            if col_lower in skip_names or any(col_lower.startswith(p)
                                              for p in ("id_", "resp_")):
                continue

            # Try numeric conversion
            s = pd.to_numeric(df[col], errors="coerce")
            non_null_ratio = s.notna().sum() / n
            if non_null_ratio < 0.85:
                continue

            # Examine unique values
            unique_vals = s.dropna().unique()
            if len(unique_vals) < 3 or len(unique_vals) > 10:
                continue

            # Skip binary check-all columns (0/1 with mostly zeros) — these
            # are usually "did you visit store X" type, not attitudinal
            if len(unique_vals) == 2:
                continue

            # Skip year-like ranges
            vmin, vmax = s.min(), s.max()
            if vmin >= 1900 and vmax <= 2100:
                continue
            # Skip very large-spread numeric columns (spend, age, etc.)
            if vmax - vmin > 20:
                continue

            candidates.append(col)

        self.log(f"Auto-detected {len(candidates)} candidate segmentation variables",
                 "info")
        return candidates

    # ---------------------------------------------------------------
    def detect_demographics(self):
        """Find demographic columns for profiling (age, income, gender, etc.).

        Uses word-boundary matching on the last segment of a column name so
        survey items like 'work in the Education industry' don't get
        mistaken for the respondent's own education.
        """
        import re
        if self.df is None:
            return {}

        # Take only the short 'tail' of each column after the last '::' or ':'
        # since survey systems prefix every sub-item with the parent question.
        def tail(col):
            s = str(col)
            for sep in (" :: ", "::", ":"):
                if sep in s:
                    s = s.rsplit(sep, 1)[-1]
            return s.strip().lower()

        # Short, specific column names (no multi-part questions) are the
        # only ones that should count as demographics.
        def is_short_column(col):
            return len(str(col)) < 120 and "::" not in str(col)

        patterns = {
            "age": r"\b(age|year born|birth year)\b",
            "gender": r"\bgender\b",
            "income": r"\b(income|household income)\b",
            "marital": r"\bmarital\b",
            "education": r"\b(education|highest level)\b",
            "employment": r"\bemployment\b",
            "ethnicity": r"\b(ethnicity|race)\b",
            "birth_year": r"\b(year born|birth year)\b",
        }

        demos = {}
        for col in self.df.columns:
            if not is_short_column(col):
                continue
            t = tail(col)
            full = str(col).lower()
            for key, pat in patterns.items():
                if key in demos:
                    continue
                if re.search(pat, t) or re.search(pat, full):
                    demos[key] = col
                    break

        self.demographics = demos
        return demos

    # ---------------------------------------------------------------
    def clean_data(self, seg_cols):
        """Prepare segmentation matrix: impute missing, convert to numeric."""
        self.log(f"Cleaning data using {len(seg_cols)} variables...", "info")
        X = self.df[seg_cols].copy()

        # Convert everything to numeric
        for c in X.columns:
            X[c] = pd.to_numeric(X[c], errors="coerce")

        # Drop rows that are ALL missing across the basis
        before = len(X)
        all_missing_mask = X.isna().all(axis=1)
        if all_missing_mask.any():
            X = X[~all_missing_mask]
            self.log(f"Dropped {before - len(X)} rows with no segmentation data",
                     "warn")

        # Impute remaining missing with column median
        n_missing = X.isna().sum().sum()
        if n_missing:
            self.log(f"Imputing {n_missing} missing cells with column medians",
                     "info")
            X = X.fillna(X.median(numeric_only=True))

        # Drop columns that have zero variance (useless for clustering)
        zero_var = X.columns[X.std() == 0].tolist()
        if zero_var:
            self.log(f"Dropping {len(zero_var)} zero-variance columns", "warn")
            X = X.drop(columns=zero_var)

        self.df_clean = X
        self.seg_cols = list(X.columns)
        self.log(f"Clean matrix: {len(X):,} rows × {len(X.columns)} variables",
                 "success")
        return X

    # ---------------------------------------------------------------
    def run_clustering(self, k=None, min_k=MIN_K, max_k=MAX_K):
        """Run K-Means clustering. If k is None, auto-select via silhouette."""
        from sklearn.preprocessing import StandardScaler
        from sklearn.cluster import KMeans
        from sklearn.metrics import silhouette_score

        X = self.df_clean
        scaler = StandardScaler()
        Xs = scaler.fit_transform(X)

        if k is None:
            self.log(f"Evaluating k from {min_k} to {max_k}...", "info")
            self.silhouette_scores = {}
            for ktest in range(min_k, max_k + 1):
                km = KMeans(n_clusters=ktest, random_state=RANDOM_STATE, n_init=20)
                lab = km.fit_predict(Xs)
                sil = silhouette_score(
                    Xs, lab,
                    sample_size=min(500, len(Xs)),
                    random_state=RANDOM_STATE
                )
                self.silhouette_scores[ktest] = sil
                self.log(f"  k={ktest}: silhouette = {sil:.4f}", "info")

            # Pick best silhouette; prefer 4-5 if within 0.01 of best for
            # actionability
            best_k = max(self.silhouette_scores, key=self.silhouette_scores.get)
            best_score = self.silhouette_scores[best_k]
            preferred = None
            for candidate in (4, 5):
                if candidate in self.silhouette_scores and \
                   best_score - self.silhouette_scores[candidate] <= 0.01:
                    preferred = candidate
                    break
            k = preferred or best_k
            self.log(f"Selected k = {k}", "success")

        self.k = k
        km = KMeans(n_clusters=k, random_state=RANDOM_STATE, n_init=50)
        labels = km.fit_predict(Xs)
        self.labels = labels + 1  # 1-indexed for presentation

        # Build profiles
        X_with_seg = X.copy()
        X_with_seg["_Segment"] = self.labels
        self.seg_profiles = X_with_seg.groupby("_Segment")[self.seg_cols].mean()
        overall_mean = X.mean()
        overall_std = X.std().replace(0, 1)  # avoid div-by-zero
        self.seg_std_profile = (self.seg_profiles - overall_mean) / overall_std

        # Auto-name segments
        self._name_segments()
        return self.labels

    # ---------------------------------------------------------------
    def _name_segments(self):
        """Assign descriptive names based on top-deviating attitudes."""
        self.seg_names = {}
        self.seg_descriptions = {}
        used = set()

        # Keyword-to-label rules (detects common market-research archetypes)
        archetypes = [
            # (keyword_in_col, direction, name, description)
            # direction: -1 means agreement with first statement of pair/item
            (["outlet", "regular"], -1, "Outlet Devotees",
             "Regularly shop outlet and factory stores; see outlets as their "
             "primary wardrobe source and a smart way to save money."),
            (["trendy", "runway"], -1, "Fashion-Forward Trendsetters",
             "First to try runway trends; view clothing as central to their "
             "identity; actively follow fashion media and celebrity style."),
            (["quality fabrics", "higher quality"], -1, "Premium Quality Seekers",
             "Invest in higher-quality fabrics (wool, cashmere, linen); willing "
             "to pay more for craftsmanship and durability."),
            (["comfort"], -1, "Comfort-First Practicals",
             "Prioritize comfort over style; low involvement with trends; "
             "practical shoppers who buy clothing they can wear daily."),
            (["bargain", "deals", "coupons"], -1, "Value-Conscious Deal-Seekers",
             "Actively hunt for deals, coupons, and sales; price-sensitive; "
             "compare prices and wait for promotions."),
            (["loyal"], -1, "Brand-Loyal Classics",
             "Loyal to trusted brands and stores; prefer proven styles and "
             "consistent quality over experimentation."),
            (["only buy", "need"], -1, "Minimalist Needs-Based",
             "Buy clothes only when needed; replace rather than refresh; "
             "low overall apparel engagement."),
            (["blend", "crowd"], 1, "Individualist Stylists",
             "Want clothing that stands out; avoid mass styles; express "
             "personal identity through unique pieces."),
        ]

        for seg in sorted(pd.Series(self.labels).unique()):
            z = self.seg_std_profile.loc[seg]
            # Find strongest-deviating attitude
            best_match = None
            best_strength = 0.3  # minimum threshold

            for keywords, direction, name, desc in archetypes:
                if name in used:
                    continue
                # find columns matching any of the keywords
                matching = [c for c in z.index
                            if all(k in str(c).lower() for k in keywords)]
                if not matching:
                    continue
                # get the most-deviated matching column in the expected direction
                for c in matching:
                    val = z[c]
                    if direction == -1 and val < -best_strength:
                        if abs(val) > best_strength:
                            best_strength = abs(val)
                            best_match = (name, desc)
                    elif direction == 1 and val > best_strength:
                        if abs(val) > best_strength:
                            best_strength = abs(val)
                            best_match = (name, desc)

            if best_match:
                self.seg_names[seg] = best_match[0]
                self.seg_descriptions[seg] = best_match[1]
                used.add(best_match[0])
            else:
                # Fallback — name by top deviating item
                top_item = z.abs().idxmax()
                top_val = z[top_item]
                direction = "low" if top_val < 0 else "high"
                name = f"Segment {seg}"
                desc = (f"Distinguished by {direction} agreement on: "
                        f"{str(top_item)[:120]}"
                        f"{'...' if len(str(top_item)) > 120 else ''}")
                self.seg_names[seg] = name
                self.seg_descriptions[seg] = desc

    # ---------------------------------------------------------------
    def segment_sizes(self):
        """Return dict of {segment_id: count}."""
        return dict(pd.Series(self.labels).value_counts().sort_index())


# ============================================================================
# REPORT GENERATION
# ============================================================================
class ReportGenerator:
    """Generates Excel workbook and PowerPoint deck from a completed
    SegmentationEngine."""

    def __init__(self, engine, log_callback=None):
        self.engine = engine
        self.log = log_callback or (lambda msg, level="info": None)

    # ---------------------------------------------------------------
    def write_excel(self, output_path):
        """Write a professional 7-sheet Excel report."""
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter
        from openpyxl.formatting.rule import ColorScaleRule

        e = self.engine
        df = e.df
        labels = e.labels

        # --- styles ---
        HEADER_FILL = PatternFill(start_color="1F4E79", end_color="1F4E79",
                                  fill_type="solid")
        HEADER_FONT = Font(name="Arial", size=11, bold=True, color="FFFFFF")
        TITLE_FONT = Font(name="Arial", size=14, bold=True, color="1F4E79")
        BIG_TITLE_FONT = Font(name="Arial", size=16, bold=True, color="1F4E79")
        TOTAL_FILL = PatternFill(start_color="D9E1F2", end_color="D9E1F2",
                                 fill_type="solid")
        BODY_FONT = Font(name="Arial", size=10)
        BOLD_BODY = Font(name="Arial", size=10, bold=True)
        CENTER = Alignment(horizontal="center", vertical="center", wrap_text=True)
        LEFT = Alignment(horizontal="left", vertical="center", wrap_text=True)
        TOP_LEFT = Alignment(horizontal="left", vertical="top", wrap_text=True)
        THIN = Side(border_style="thin", color="BFBFBF")
        BORDER = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)

        def write_table(ws, data_df, start_row=1, title=None,
                        pct_cols=None, num_cols=None, dollar_cols=None,
                        freeze_first_col=True):
            if title:
                ws.cell(row=start_row, column=1, value=title).font = TITLE_FONT
                hdr_row = start_row + 2
            else:
                hdr_row = start_row

            for j, c in enumerate(data_df.columns, 1):
                cell = ws.cell(row=hdr_row, column=j, value=str(c))
                cell.font = HEADER_FONT
                cell.fill = HEADER_FILL
                cell.alignment = CENTER
                cell.border = BORDER

            for i, (_, r) in enumerate(data_df.iterrows(), hdr_row + 1):
                first_val = str(r.iloc[0]).upper() if len(r) else ""
                is_total = first_val in ("TOTAL", "OVERALL", "ALL")
                for j, col in enumerate(data_df.columns, 1):
                    v = r[col]
                    cell = ws.cell(row=i, column=j,
                                   value=("" if pd.isna(v) else v))
                    cell.font = BOLD_BODY if is_total else BODY_FONT
                    cell.border = BORDER
                    if isinstance(v, (int, float, np.integer, np.floating)) and j > 1:
                        cell.alignment = CENTER
                        if pct_cols and col in pct_cols:
                            cell.number_format = '0.0"%"'
                        elif dollar_cols and col in dollar_cols:
                            cell.number_format = "$#,##0"
                        elif num_cols and col in num_cols:
                            cell.number_format = "0.00"
                    else:
                        cell.alignment = LEFT if j <= 2 else CENTER
                    if is_total:
                        cell.fill = TOTAL_FILL

            # column widths
            for j, col in enumerate(data_df.columns, 1):
                sample = [str(col)] + [str(r[col])[:50]
                                        for _, r in data_df.iterrows()]
                w = min(max(max(len(s) for s in sample) + 3, 12), 45)
                ws.column_dimensions[get_column_letter(j)].width = w

            ws.row_dimensions[hdr_row].height = 32
            return hdr_row + len(data_df)  # last filled row

        wb = Workbook()

        # ----------------------------------------------------------
        # Sheet 1: Executive Summary
        # ----------------------------------------------------------
        ws = wb.active
        ws.title = "Executive Summary"
        ws.cell(row=1, column=1,
                value=f"Customer Segmentation Report").font = BIG_TITLE_FONT
        subtitle = (f"{e.k} segments identified from {len(df):,} respondents "
                    f"using {len(e.seg_cols)} variables "
                    f"· Generated {datetime.now():%Y-%m-%d %H:%M}")
        ws.cell(row=2, column=1, value=subtitle).font = Font(
            name="Arial", size=10, italic=True, color="404040")

        sizes = e.segment_sizes()
        total_n = sum(sizes.values())
        summary_rows = []
        for seg, size in sorted(sizes.items()):
            summary_rows.append({
                "Segment": seg,
                "Segment Name": e.seg_names.get(seg, f"Segment {seg}"),
                "Size (n)": size,
                "Size (%)": size / total_n * 100,
                "Description": e.seg_descriptions.get(seg, ""),
            })
        summary_df = pd.DataFrame(summary_rows)
        write_table(ws, summary_df, start_row=4,
                    title="Segment Overview",
                    pct_cols=["Size (%)"])

        # ----------------------------------------------------------
        # Sheet 2: Demographics
        # ----------------------------------------------------------
        ws = wb.create_sheet("Demographics")
        demo_df = self._build_demographics_table()
        if demo_df is not None and len(demo_df):
            pct_cols = [c for c in demo_df.columns if c.startswith("%")]
            num_cols = [c for c in demo_df.columns
                        if c.startswith("Avg") or "Mean" in c]
            write_table(ws, demo_df, start_row=1,
                        title="Segment Demographics",
                        pct_cols=pct_cols + ["Size (%)"],
                        num_cols=num_cols)
        else:
            ws.cell(row=1, column=1,
                    value="No demographic columns were detected in the input file."
                    ).font = TITLE_FONT

        # ----------------------------------------------------------
        # Sheet 3: Segment Profiles (all basis variables)
        # ----------------------------------------------------------
        ws = wb.create_sheet("Segment Profiles")
        prof_df = e.seg_profiles.T.reset_index()
        prof_df.columns = ["Variable"] + [
            f"Seg {s}: {e.seg_names.get(s, f'Segment {s}')}"
            for s in e.seg_profiles.index
        ]
        prof_df["Overall"] = [e.df_clean[v].mean() for v in prof_df["Variable"]]
        num_cols = [c for c in prof_df.columns if c != "Variable"]
        write_table(ws, prof_df, start_row=1,
                    title="Mean Score on Each Variable by Segment",
                    num_cols=num_cols)

        # ----------------------------------------------------------
        # Sheet 4: Top Differentiators
        # ----------------------------------------------------------
        ws = wb.create_sheet("Key Differentiators")
        rows = []
        for seg in e.seg_std_profile.index:
            z = e.seg_std_profile.loc[seg]
            top = z.abs().sort_values(ascending=False).head(10)
            for rank, (var, _) in enumerate(top.items(), 1):
                rows.append({
                    "Segment": seg,
                    "Segment Name": e.seg_names.get(seg, f"Segment {seg}"),
                    "Rank": rank,
                    "Variable": var,
                    "Segment Mean": e.seg_profiles.loc[seg, var],
                    "Overall Mean": e.df_clean[var].mean(),
                    "Z-Score": z[var],
                })
        diff_df = pd.DataFrame(rows)
        write_table(ws, diff_df, start_row=1,
                    title="Top 10 Most-Distinguishing Variables per Segment",
                    num_cols=["Segment Mean", "Overall Mean", "Z-Score"])

        # ----------------------------------------------------------
        # Sheet 5: Z-Score Heatmap
        # ----------------------------------------------------------
        ws = wb.create_sheet("Z-Score Heatmap")
        z_df = e.seg_std_profile.T.reset_index()
        z_df.columns = ["Variable"] + [
            f"Seg {s}: {e.seg_names.get(s, f'Segment {s}')}"
            for s in e.seg_std_profile.index
        ]
        last_row = write_table(
            ws, z_df, start_row=1,
            title="How Much Each Segment Deviates from the Overall Mean "
                  "(in standard deviations)",
            num_cols=[c for c in z_df.columns if c != "Variable"])

        first_col = 2
        last_col = len(z_df.columns)
        rng = (f"{get_column_letter(first_col)}4:"
               f"{get_column_letter(last_col)}{last_row}")
        ws.conditional_formatting.add(rng, ColorScaleRule(
            start_type="num", start_value=-2, start_color="C5504B",
            mid_type="num", mid_value=0, mid_color="FFFFFF",
            end_type="num", end_value=2, end_color="0F7B3C"
        ))

        # ----------------------------------------------------------
        # Sheet 6: Respondent Assignments
        # ----------------------------------------------------------
        ws = wb.create_sheet("Respondent Assignments")
        assign_df = pd.DataFrame({
            "Respondent #": range(1, len(labels) + 1),
            "Segment": labels,
            "Segment Name": [e.seg_names.get(s, f"Segment {s}") for s in labels],
        })
        # Attach any detected demographics for reference
        demos = e.demographics
        for key, col in demos.items():
            if col in df.columns and len(df) == len(labels):
                assign_df[key.title()] = df[col].values
        write_table(ws, assign_df, start_row=1,
                    title="Per-Respondent Segment Assignment")

        # ----------------------------------------------------------
        # Sheet 7: Methodology
        # ----------------------------------------------------------
        ws = wb.create_sheet("Methodology")
        ws.column_dimensions["A"].width = 110
        lines = [
            ("TITLE", "Segmentation Methodology"),
            ("", ""),
            ("H2", "INPUT"),
            ("P", f"• {len(df):,} records loaded from the input file"),
            ("P", f"• {len(e.seg_cols)} variables used as segmentation basis"),
            ("P", f"• Basis variables selected by: "
                  f"auto-detection of Likert-style items with 85%+ completeness"),
            ("", ""),
            ("H2", "DATA CLEANING"),
            ("P", "• Fully-empty rows and columns dropped"),
            ("P", "• Missing values imputed with column median"),
            ("P", "• Zero-variance columns removed"),
            ("P", "• All variables converted to numeric"),
            ("", ""),
            ("H2", "CLUSTERING ALGORITHM"),
            ("P", "• Standardization: z-score (mean=0, std=1) across all basis "
                  "variables"),
            ("P", "• Method: K-Means clustering with 50 random initializations"),
            ("P", "• Random seed: 42 (results are fully reproducible)"),
            ("P", f"• k evaluation range: {MIN_K} to {MAX_K}"),
            ("P", f"• Selected k = {e.k}"),
            ("P", f"• Silhouette scores by k: "
                  f"{', '.join(f'k={k}: {s:.3f}' for k, s in e.silhouette_scores.items())}"),
            ("", ""),
            ("H2", "SEGMENT NAMING"),
            ("P", "• Each segment auto-labeled using its top-deviating variables"),
            ("P", "• Naming rules match common market-research archetypes: "
                  "outlet devotees, trendsetters, quality seekers, comfort-first, "
                  "deal-seekers, brand loyalists, minimalists, individualists"),
            ("P", "• If no archetype matches, segment is labeled by its "
                  "strongest-deviating variable"),
            ("", ""),
            ("H2", "HOW TO READ THE OUTPUT"),
            ("P", "• Z-Score = (Segment Mean − Overall Mean) ÷ Overall Std. Dev."),
            ("P", "• |Z| > 0.5 = meaningful deviation; |Z| > 1.0 = strong deviation"),
            ("P", "• Green cells in the Z-Score Heatmap = segment scores higher "
                  "than average"),
            ("P", "• Red cells in the Z-Score Heatmap = segment scores lower "
                  "than average"),
            ("", ""),
            ("H2", f"REPORT GENERATED"),
            ("P", f"• Generated by {APP_NAME} v{APP_VERSION}"),
            ("P", f"• Timestamp: {datetime.now():%Y-%m-%d %H:%M:%S}"),
        ]
        for i, (kind, text) in enumerate(lines, 1):
            cell = ws.cell(row=i, column=1, value=text)
            if kind == "TITLE":
                cell.font = BIG_TITLE_FONT
            elif kind == "H2":
                cell.font = Font(name="Arial", size=12, bold=True,
                                 color="1F4E79")
            else:
                cell.font = BODY_FONT
            cell.alignment = TOP_LEFT

        # Freeze panes on main sheets
        for sht in ("Executive Summary", "Demographics", "Segment Profiles",
                    "Key Differentiators", "Z-Score Heatmap",
                    "Respondent Assignments"):
            if sht in wb.sheetnames:
                wb[sht].freeze_panes = "B5" if sht != "Respondent Assignments" \
                                       else "B2"

        wb.save(output_path)
        self.log(f"✓ Excel report saved: {output_path}", "success")

    # ---------------------------------------------------------------
    def _build_demographics_table(self):
        """Build a demographics-by-segment table using detected demo columns."""
        e = self.engine
        df = e.df.reset_index(drop=True)
        if len(df) != len(e.labels):
            # Data cleaning dropped rows; align
            df = df.iloc[:len(e.labels)].copy()
        df = df.copy()
        df["_Segment"] = e.labels
        demos = e.demographics

        rows = []
        all_rows = list(range(len(df)))

        def pct(mask):
            return mask.sum() / len(mask) * 100 if len(mask) else 0

        for seg in sorted(df["_Segment"].unique()):
            sub = df[df["_Segment"] == seg]
            size = len(sub)
            row = {
                "Segment": seg,
                "Segment Name": e.seg_names.get(seg, f"Segment {seg}"),
                "Size (n)": size,
                "Size (%)": size / len(df) * 100,
            }
            # Age
            if "age" in demos and demos["age"] in df.columns:
                vals = pd.to_numeric(sub[demos["age"]], errors="coerce")
                row["Avg Age"] = vals.mean()
            elif "birth_year" in demos and demos["birth_year"] in df.columns:
                vals = pd.to_numeric(sub[demos["birth_year"]], errors="coerce")
                # If encoded (e.g. 1-80 maps to 1931-2010 birth years), treat
                # as a raw proxy for age
                row["Avg (Year Born encoded)"] = vals.mean()
            # Income
            if "income" in demos and demos["income"] in df.columns:
                col = demos["income"]
                vals = sub[col]
                # If numeric and binary 1/2, label it
                nunique = vals.nunique()
                if nunique <= 5:
                    for v in sorted(vals.dropna().unique()):
                        row[f"% Income={v}"] = pct(vals == v)
                else:
                    vals = pd.to_numeric(vals, errors="coerce")
                    row["Avg Income"] = vals.mean()
            # Other categorical demographics
            for key in ("marital", "education", "employment", "ethnicity",
                        "gender"):
                if key in demos and demos[key] in df.columns:
                    col = demos[key]
                    if df[col].nunique() <= 5:
                        # Most common value share
                        top_val = df[col].mode().iloc[0] \
                            if not df[col].mode().empty else None
                        if top_val is not None:
                            row[f"% {key.title()}={top_val}"] = pct(
                                sub[col] == top_val)
            rows.append(row)

        # Overall row
        all_row = {
            "Segment": "TOTAL",
            "Segment Name": "All Respondents",
            "Size (n)": len(df),
            "Size (%)": 100.0,
        }
        for r in rows:
            for k in r:
                if k in ("Segment", "Segment Name", "Size (n)", "Size (%)"):
                    continue
                if k not in all_row:
                    # Compute overall
                    if k.startswith("Avg Age") and "age" in demos:
                        all_row[k] = pd.to_numeric(
                            df[demos["age"]], errors="coerce").mean()
                    elif k.startswith("Avg Income") and "income" in demos:
                        all_row[k] = pd.to_numeric(
                            df[demos["income"]], errors="coerce").mean()
                    elif k.startswith("Avg (Year Born encoded)") and "birth_year" in demos:
                        all_row[k] = pd.to_numeric(
                            df[demos["birth_year"]], errors="coerce").mean()
                    elif k.startswith("% Income="):
                        v = k.split("=")[1]
                        try:
                            v_num = float(v)
                            all_row[k] = pct(
                                pd.to_numeric(df[demos["income"]],
                                              errors="coerce") == v_num)
                        except ValueError:
                            all_row[k] = pct(df[demos["income"]] == v)
                    elif k.startswith("%"):
                        # Other categorical
                        key = k.split(" ")[1].split("=")[0].lower()
                        val = k.split("=")[1]
                        if key in demos:
                            col = demos[key]
                            all_row[k] = pct(df[col].astype(str) == val)
        rows.append(all_row)
        return pd.DataFrame(rows)

    # ---------------------------------------------------------------
    def write_pptx(self, output_path):
        """Write a PowerPoint deck summarizing the segmentation."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dgm.color import RGBColor  # wrong import on purpose to trigger except below if needed
        except ImportError:
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt, Emu
                from pptx.dml.color import RGBColor
                from pptx.enum.shapes import MSO_SHAPE
                from pptx.enum.text import PP_ALIGN
            except ImportError as e:
                self.log(f"python-pptx not installed — skipping deck: {e}", "error")
                raise
        else:
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.enum.text import PP_ALIGN

        e = self.engine
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # ---- Color palette: Midnight Executive ----
        NAVY = RGBColor(0x1F, 0x4E, 0x79)
        NAVY_DARK = RGBColor(0x16, 0x3B, 0x5C)
        ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
        TEXT_GREY = RGBColor(0x5A, 0x67, 0x73)
        ACCENT = RGBColor(0xE8, 0x8B, 0x3A)   # warm contrast

        blank = prs.slide_layouts[6]

        def add_rect(slide, x, y, w, h, fill, line=None):
            shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
            if line is None:
                shp.line.fill.background()
            else:
                shp.line.color.rgb = line
            shp.shadow.inherit = False
            return shp

        def add_text(slide, x, y, w, h, text, font_size=14,
                     color=TEXT_DARK, bold=False, align=PP_ALIGN.LEFT,
                     font_name="Calibri"):
            box = slide.shapes.add_textbox(x, y, w, h)
            tf = box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = Inches(0.05)
            tf.margin_top = tf.margin_bottom = Inches(0.02)
            p = tf.paragraphs[0]
            p.alignment = align
            r = p.add_run()
            r.text = text
            r.font.size = Pt(font_size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font_name
            return box

        # ============================================================
        # Slide 1 — Title
        # ============================================================
        s = prs.slides.add_slide(blank)
        add_rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
        # Decorative accent square
        add_rect(s, Inches(0.75), Inches(3.0), Inches(0.4), Inches(0.4), ACCENT)
        add_text(s, Inches(0.75), Inches(3.5), Inches(12), Inches(1.2),
                 "Customer Segmentation",
                 font_size=54, color=WHITE, bold=True, font_name="Georgia")
        add_text(s, Inches(0.75), Inches(4.5), Inches(12), Inches(0.6),
                 f"{e.k} segments · {len(e.df):,} respondents · "
                 f"{len(e.seg_cols)} variables",
                 font_size=20, color=ICE_BLUE, font_name="Calibri")
        add_text(s, Inches(0.75), Inches(6.7), Inches(12), Inches(0.4),
                 f"Generated by {APP_NAME} · {datetime.now():%B %d, %Y}",
                 font_size=12, color=ICE_BLUE, font_name="Calibri")

        # ============================================================
        # Slide 2 — Segment Overview (cards)
        # ============================================================
        s = prs.slides.add_slide(blank)
        add_text(s, Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.7),
                 "Segment Overview",
                 font_size=32, color=NAVY, bold=True, font_name="Georgia")
        add_text(s, Inches(0.5), Inches(1.05), Inches(12.5), Inches(0.4),
                 f"{e.k} distinct customer segments identified",
                 font_size=14, color=TEXT_GREY, font_name="Calibri")

        sizes = e.segment_sizes()
        total = sum(sizes.values())
        segs = sorted(sizes.keys())
        n_segs = len(segs)
        # Layout cards in a row (up to 4 across)
        cols = min(n_segs, 4)
        rows_needed = (n_segs + cols - 1) // cols
        card_w = Inches((12.3) / cols - 0.15)
        card_h = Inches(2.5 if rows_needed == 1 else 2.3)
        start_y = Inches(1.7)
        gap_x = Inches(0.15)
        gap_y = Inches(0.25)

        for i, seg in enumerate(segs):
            col = i % cols
            row = i // cols
            x = Inches(0.5) + (card_w + gap_x) * col
            y = start_y + (card_h + gap_y) * row
            # Card background
            add_rect(s, x, y, card_w, card_h, WHITE, line=ICE_BLUE)
            # Accent strip on top
            add_rect(s, x, y, card_w, Inches(0.12), NAVY)
            # Segment number badge
            add_rect(s, x + Inches(0.25), y + Inches(0.3),
                     Inches(0.5), Inches(0.5), ACCENT)
            add_text(s, x + Inches(0.25), y + Inches(0.33),
                     Inches(0.5), Inches(0.5),
                     str(seg), font_size=22, color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER, font_name="Georgia")
            # Name
            name = e.seg_names.get(seg, f"Segment {seg}")
            add_text(s, x + Inches(0.9), y + Inches(0.3),
                     card_w - Inches(1.0), Inches(0.6),
                     name, font_size=15, color=NAVY, bold=True,
                     font_name="Calibri")
            # Big size stat
            pct = sizes[seg] / total * 100
            add_text(s, x + Inches(0.25), y + Inches(0.95),
                     card_w - Inches(0.5), Inches(0.7),
                     f"{pct:.0f}%", font_size=40, color=NAVY, bold=True,
                     font_name="Georgia")
            add_text(s, x + Inches(0.25), y + Inches(1.6),
                     card_w - Inches(0.5), Inches(0.3),
                     f"n = {sizes[seg]:,}",
                     font_size=11, color=TEXT_GREY, font_name="Calibri")
            # Description (trimmed)
            desc = e.seg_descriptions.get(seg, "")
            if len(desc) > 140:
                desc = desc[:137].rsplit(" ", 1)[0] + "..."
            add_text(s, x + Inches(0.25), y + Inches(1.9),
                     card_w - Inches(0.5), card_h - Inches(2.0),
                     desc, font_size=10, color=TEXT_DARK,
                     font_name="Calibri")

        # ============================================================
        # One slide per segment — deep-dive
        # ============================================================
        for seg in segs:
            s = prs.slides.add_slide(blank)
            name = e.seg_names.get(seg, f"Segment {seg}")
            desc = e.seg_descriptions.get(seg, "")

            # Left panel — segment header (navy)
            add_rect(s, 0, 0, Inches(4.5), prs.slide_height, NAVY)
            add_rect(s, Inches(0.5), Inches(0.5),
                     Inches(0.7), Inches(0.7), ACCENT)
            add_text(s, Inches(0.5), Inches(0.55),
                     Inches(0.7), Inches(0.6),
                     str(seg), font_size=32, color=WHITE, bold=True,
                     align=PP_ALIGN.CENTER, font_name="Georgia")
            add_text(s, Inches(0.5), Inches(1.5),
                     Inches(3.5), Inches(0.5),
                     "SEGMENT", font_size=11, color=ICE_BLUE,
                     font_name="Calibri", bold=True)
            add_text(s, Inches(0.5), Inches(1.8),
                     Inches(3.6), Inches(1.6),
                     name, font_size=26, color=WHITE, bold=True,
                     font_name="Georgia")
            # Big stat — size
            pct = sizes[seg] / total * 100
            add_text(s, Inches(0.5), Inches(4.8),
                     Inches(3.6), Inches(1.0),
                     f"{pct:.0f}%", font_size=64, color=WHITE, bold=True,
                     font_name="Georgia")
            add_text(s, Inches(0.5), Inches(6.0),
                     Inches(3.6), Inches(0.4),
                     f"{sizes[seg]:,} respondents",
                     font_size=14, color=ICE_BLUE, font_name="Calibri")

            # Right panel — content
            # Description
            add_text(s, Inches(4.9), Inches(0.5),
                     Inches(8.0), Inches(0.4),
                     "WHO THEY ARE", font_size=11, color=NAVY,
                     bold=True, font_name="Calibri")
            add_text(s, Inches(4.9), Inches(0.85),
                     Inches(8.0), Inches(1.4),
                     desc, font_size=15, color=TEXT_DARK,
                     font_name="Calibri")

            # Top distinguishing attitudes
            add_text(s, Inches(4.9), Inches(2.4),
                     Inches(8.0), Inches(0.4),
                     "KEY DIFFERENTIATORS", font_size=11, color=NAVY,
                     bold=True, font_name="Calibri")
            z = e.seg_std_profile.loc[seg]
            top5 = z.abs().sort_values(ascending=False).head(5)
            y_pos = Inches(2.8)
            for i, (var, _) in enumerate(top5.items()):
                direction = "↑ higher" if z[var] > 0 else "↓ lower"
                color = RGBColor(0x0F, 0x7B, 0x3C) if z[var] > 0 \
                    else RGBColor(0xC5, 0x50, 0x4B)
                # Bullet dot
                add_rect(s, Inches(4.95), y_pos + Inches(0.12),
                         Inches(0.08), Inches(0.08), NAVY)
                # Variable name
                var_short = smart_short(var, max_len=95)
                add_text(s, Inches(5.15), y_pos,
                         Inches(6.2), Inches(0.7),
                         var_short, font_size=11, color=TEXT_DARK,
                         font_name="Calibri")
                # Direction badge
                add_text(s, Inches(11.4), y_pos,
                         Inches(1.5), Inches(0.4),
                         f"{direction} ({z[var]:+.2f}σ)",
                         font_size=11, color=color, bold=True,
                         font_name="Calibri")
                y_pos += Inches(0.75)

        # ============================================================
        # Methodology slide
        # ============================================================
        s = prs.slides.add_slide(blank)
        add_text(s, Inches(0.5), Inches(0.3),
                 Inches(12.5), Inches(0.7),
                 "Methodology",
                 font_size=32, color=NAVY, bold=True, font_name="Georgia")
        add_text(s, Inches(0.5), Inches(1.05),
                 Inches(12.5), Inches(0.4),
                 "How these segments were built",
                 font_size=14, color=TEXT_GREY, font_name="Calibri")

        # Two-column layout
        col_w = Inches(6.0)
        col1_x = Inches(0.5)
        col2_x = Inches(6.8)
        col_y = Inches(1.8)
        col_h = Inches(5.5)

        # Left column — Process
        add_rect(s, col1_x, col_y, col_w, col_h, WHITE, line=ICE_BLUE)
        add_rect(s, col1_x, col_y, col_w, Inches(0.5), NAVY)
        add_text(s, col1_x + Inches(0.2), col_y + Inches(0.08),
                 col_w, Inches(0.35),
                 "PROCESS", font_size=14, color=WHITE, bold=True,
                 font_name="Calibri")

        steps = [
            ("01", "Load & inspect",
             f"{len(e.df):,} records from the uploaded file"),
            ("02", "Auto-detect basis",
             f"{len(e.seg_cols)} Likert-style attitudinal items"),
            ("03", "Clean data",
             "Drop empty rows, impute missing with median, remove zero-var"),
            ("04", "Standardize & cluster",
             "Z-score normalize, then K-Means (50 inits, seed=42)"),
            ("05", "Select k",
             f"Silhouette tested k={MIN_K}–{MAX_K}, chose k={e.k}"),
            ("06", "Name & profile",
             "Auto-label from top-deviating attitudes"),
        ]
        y = col_y + Inches(0.7)
        for num, title, desc in steps:
            add_text(s, col1_x + Inches(0.3), y,
                     Inches(0.5), Inches(0.4),
                     num, font_size=16, color=ACCENT, bold=True,
                     font_name="Georgia")
            add_text(s, col1_x + Inches(0.9), y - Inches(0.02),
                     col_w - Inches(1.1), Inches(0.35),
                     title, font_size=13, color=NAVY, bold=True,
                     font_name="Calibri")
            add_text(s, col1_x + Inches(0.9), y + Inches(0.3),
                     col_w - Inches(1.1), Inches(0.45),
                     desc, font_size=10, color=TEXT_GREY,
                     font_name="Calibri")
            y += Inches(0.78)

        # Right column — Silhouette / quality
        add_rect(s, col2_x, col_y, col_w, col_h, WHITE, line=ICE_BLUE)
        add_rect(s, col2_x, col_y, col_w, Inches(0.5), NAVY)
        add_text(s, col2_x + Inches(0.2), col_y + Inches(0.08),
                 col_w, Inches(0.35),
                 "CLUSTER QUALITY", font_size=14, color=WHITE, bold=True,
                 font_name="Calibri")

        add_text(s, col2_x + Inches(0.3), col_y + Inches(0.7),
                 col_w - Inches(0.6), Inches(0.4),
                 "Silhouette score by k",
                 font_size=12, color=TEXT_DARK, bold=True,
                 font_name="Calibri")
        add_text(s, col2_x + Inches(0.3), col_y + Inches(1.05),
                 col_w - Inches(0.6), Inches(0.5),
                 "Higher = cleaner separation. Max value is 1.0.",
                 font_size=10, color=TEXT_GREY, font_name="Calibri")

        # Silhouette bars
        scores = e.silhouette_scores
        if scores:
            bar_y = col_y + Inches(1.7)
            max_sil = max(scores.values()) if scores else 1
            max_bar_w = col_w - Inches(2.0)
            for kk, sil in sorted(scores.items()):
                is_selected = (kk == e.k)
                # k label
                add_text(s, col2_x + Inches(0.3), bar_y,
                         Inches(0.6), Inches(0.3),
                         f"k={kk}", font_size=11,
                         color=(NAVY if is_selected else TEXT_DARK),
                         bold=is_selected, font_name="Calibri")
                # Bar
                bar_w = Emu(int(max_bar_w * (sil / max(max_sil, 0.01))))
                bar_color = ACCENT if is_selected else ICE_BLUE
                add_rect(s, col2_x + Inches(0.9), bar_y + Inches(0.05),
                         bar_w, Inches(0.2), bar_color)
                # Score text
                add_text(s, col2_x + Inches(0.9) + max_bar_w + Inches(0.1),
                         bar_y,
                         Inches(0.8), Inches(0.3),
                         f"{sil:.3f}", font_size=11,
                         color=(NAVY if is_selected else TEXT_DARK),
                         bold=is_selected, font_name="Calibri")
                bar_y += Inches(0.38)

            add_text(s, col2_x + Inches(0.3), bar_y + Inches(0.2),
                     col_w - Inches(0.6), Inches(0.4),
                     f"Selected: k = {e.k}",
                     font_size=12, color=NAVY, bold=True,
                     font_name="Calibri")

        # Closing slide
        s = prs.slides.add_slide(blank)
        add_rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
        add_text(s, Inches(0.75), Inches(3.0), Inches(12), Inches(1.2),
                 "Next Steps",
                 font_size=48, color=WHITE, bold=True, font_name="Georgia")
        next_steps = (
            "• Review the Excel workbook for full segment profiles and "
            "respondent-level assignments\n"
            "• Validate segments with business stakeholders\n"
            "• Size target segments using the assignment file\n"
            "• Design segment-specific messaging and product strategy"
        )
        add_text(s, Inches(0.75), Inches(4.2), Inches(12), Inches(2.5),
                 next_steps, font_size=18, color=ICE_BLUE,
                 font_name="Calibri")

        prs.save(output_path)
        self.log(f"✓ PowerPoint deck saved: {output_path}", "success")


# ============================================================================
# TKINTER APP
# ============================================================================
if _TK_AVAILABLE:
    _APP_BASE = tk.Tk
else:
    _APP_BASE = object


class SegmentProApp(_APP_BASE):

    def __init__(self):
        super().__init__()
        self.title(f"{APP_NAME} — {APP_SUBTITLE}")
        self.geometry("1100x760")
        self.minsize(950, 640)
        self.configure(bg=COLOR_BG)

        # State
        self.engine = SegmentationEngine(log_callback=self._log)
        self.input_filepath = None
        self.output_dir = None
        self.detected_cols = []
        self.selected_cols = set()
        self.current_step = 1
        self.is_running = False

        self._setup_style()
        self._build_ui()

    # ---------------------------------------------------------------
    def _setup_style(self):
        style = ttk.Style(self)
        try:
            style.theme_use("clam")
        except tk.TclError:
            pass

        style.configure("TFrame", background=COLOR_BG)
        style.configure("Panel.TFrame", background=COLOR_PANEL,
                        relief="solid", borderwidth=1)
        style.configure("TLabel", background=COLOR_BG, foreground=COLOR_TEXT,
                        font=FONT_BODY)
        style.configure("Panel.TLabel", background=COLOR_PANEL,
                        foreground=COLOR_TEXT, font=FONT_BODY)
        style.configure("Title.TLabel", background=COLOR_BG,
                        foreground=COLOR_PRIMARY, font=FONT_TITLE)
        style.configure("Subtitle.TLabel", background=COLOR_BG,
                        foreground=COLOR_TEXT_MUTED, font=FONT_SUBTITLE)
        style.configure("H2.TLabel", background=COLOR_PANEL,
                        foreground=COLOR_PRIMARY, font=FONT_H2)
        style.configure("Muted.TLabel", background=COLOR_PANEL,
                        foreground=COLOR_TEXT_MUTED, font=FONT_SMALL)

        style.configure("Primary.TButton",
                        font=("Segoe UI", 10, "bold"),
                        foreground="#FFFFFF",
                        background=COLOR_PRIMARY,
                        padding=(16, 10),
                        borderwidth=0)
        style.map("Primary.TButton",
                  background=[("active", COLOR_PRIMARY_DARK),
                              ("disabled", "#9AB3C8")],
                  foreground=[("disabled", "#EEEEEE")])

        style.configure("Secondary.TButton",
                        font=("Segoe UI", 10),
                        padding=(14, 8),
                        borderwidth=1)

        style.configure("Success.TLabel", background=COLOR_PANEL,
                        foreground=COLOR_SUCCESS, font=FONT_SMALL)
        style.configure("Error.TLabel", background=COLOR_PANEL,
                        foreground=COLOR_ERROR, font=FONT_SMALL)

        style.configure("TNotebook", background=COLOR_BG, borderwidth=0)
        style.configure("TNotebook.Tab",
                        font=("Segoe UI", 10, "bold"),
                        padding=(18, 10))

        style.configure("Horizontal.TProgressbar",
                        background=COLOR_PRIMARY,
                        troughcolor="#E0E6EE",
                        borderwidth=0,
                        thickness=14)

        # Treeview
        style.configure("Treeview",
                        background=COLOR_PANEL,
                        foreground=COLOR_TEXT,
                        rowheight=26,
                        fieldbackground=COLOR_PANEL,
                        font=FONT_BODY,
                        borderwidth=0)
        style.configure("Treeview.Heading",
                        font=("Segoe UI", 9, "bold"),
                        background=COLOR_PRIMARY,
                        foreground="#FFFFFF",
                        relief="flat")
        style.map("Treeview.Heading",
                  background=[("active", COLOR_PRIMARY_DARK)])
        style.map("Treeview",
                  background=[("selected", COLOR_ACCENT)],
                  foreground=[("selected", "#FFFFFF")])

    # ---------------------------------------------------------------
    def _build_ui(self):
        # Header bar
        header = tk.Frame(self, bg=COLOR_PRIMARY, height=80)
        header.pack(fill="x", side="top")
        header.pack_propagate(False)

        title_frame = tk.Frame(header, bg=COLOR_PRIMARY)
        title_frame.pack(side="left", padx=24, pady=16)
        tk.Label(title_frame, text=APP_NAME,
                 font=("Segoe UI", 20, "bold"),
                 fg="#FFFFFF", bg=COLOR_PRIMARY).pack(anchor="w")
        tk.Label(title_frame, text=APP_SUBTITLE,
                 font=("Segoe UI", 10),
                 fg=COLOR_BORDER, bg=COLOR_PRIMARY).pack(anchor="w")

        version_frame = tk.Frame(header, bg=COLOR_PRIMARY)
        version_frame.pack(side="right", padx=24, pady=24)
        tk.Label(version_frame, text=f"v{APP_VERSION}",
                 font=("Segoe UI", 10),
                 fg=COLOR_BORDER, bg=COLOR_PRIMARY).pack()

        # Main notebook
        self.notebook = ttk.Notebook(self)
        self.notebook.pack(fill="both", expand=True, padx=16, pady=(14, 8))

        self.tab_load = ttk.Frame(self.notebook, style="TFrame")
        self.tab_vars = ttk.Frame(self.notebook, style="TFrame")
        self.tab_run = ttk.Frame(self.notebook, style="TFrame")
        self.tab_results = ttk.Frame(self.notebook, style="TFrame")

        self.notebook.add(self.tab_load, text="  1. Load Data  ")
        self.notebook.add(self.tab_vars, text="  2. Choose Variables  ")
        self.notebook.add(self.tab_run, text="  3. Run Analysis  ")
        self.notebook.add(self.tab_results, text="  4. Results & Export  ")

        self._build_tab_load()
        self._build_tab_vars()
        self._build_tab_run()
        self._build_tab_results()

        # Footer / log
        self._build_footer()

        # Disable tabs after the first until ready
        self._update_tab_state()

    # ---------------------------------------------------------------
    def _build_tab_load(self):
        tab = self.tab_load
        panel = ttk.Frame(tab, style="Panel.TFrame", padding=30)
        panel.pack(fill="both", expand=True, padx=16, pady=16)

        ttk.Label(panel, text="Step 1 — Load your data",
                  style="H2.TLabel").pack(anchor="w")
        ttk.Label(panel,
                  text="Upload an Excel (.xlsx, .xls) or CSV file. Each row "
                       "should be one respondent; columns are survey items.",
                  style="Muted.TLabel").pack(anchor="w", pady=(4, 20))

        # File selector box
        box = tk.Frame(panel, bg="#F0F4FA", relief="solid", bd=1,
                       highlightthickness=0)
        box.pack(fill="x", pady=10)
        inner = tk.Frame(box, bg="#F0F4FA")
        inner.pack(padx=20, pady=26)

        tk.Label(inner, text="📁", font=("Segoe UI Emoji", 32),
                 bg="#F0F4FA", fg=COLOR_PRIMARY).pack()
        self.file_label = tk.Label(
            inner, text="No file selected",
            font=("Segoe UI", 11),
            bg="#F0F4FA", fg=COLOR_TEXT_MUTED)
        self.file_label.pack(pady=(8, 14))

        ttk.Button(inner, text="Browse for file…",
                   style="Primary.TButton",
                   command=self._choose_file).pack()

        # Dataset summary (hidden until loaded)
        self.dataset_summary = tk.Frame(panel, bg=COLOR_PANEL)
        self.dataset_summary.pack(fill="x", pady=(22, 0))

        # Sample preview (hidden until loaded)
        self.preview_frame = tk.Frame(panel, bg=COLOR_PANEL)
        self.preview_frame.pack(fill="both", expand=True, pady=(14, 0))

        # Next button
        nav = tk.Frame(panel, bg=COLOR_PANEL)
        nav.pack(fill="x", pady=(14, 0))
        self.btn_to_vars = ttk.Button(
            nav, text="Next: Choose Variables →",
            style="Primary.TButton",
            command=lambda: self.notebook.select(self.tab_vars),
            state="disabled")
        self.btn_to_vars.pack(side="right")

    # ---------------------------------------------------------------
    def _build_tab_vars(self):
        tab = self.tab_vars
        panel = ttk.Frame(tab, style="Panel.TFrame", padding=30)
        panel.pack(fill="both", expand=True, padx=16, pady=16)

        ttk.Label(panel, text="Step 2 — Choose segmentation variables",
                  style="H2.TLabel").pack(anchor="w")
        ttk.Label(panel,
                  text="These are the variables used to distinguish segments. "
                       "Auto-detected Likert-style items are pre-selected. "
                       "Check / uncheck to adjust.",
                  style="Muted.TLabel").pack(anchor="w", pady=(4, 16))

        # Toolbar
        tools = tk.Frame(panel, bg=COLOR_PANEL)
        tools.pack(fill="x", pady=(0, 10))

        self.var_count_label = tk.Label(
            tools, text="0 variables selected",
            font=("Segoe UI", 10, "bold"),
            bg=COLOR_PANEL, fg=COLOR_PRIMARY)
        self.var_count_label.pack(side="left")

        ttk.Button(tools, text="Select All",
                   style="Secondary.TButton",
                   command=self._select_all_vars).pack(side="right", padx=4)
        ttk.Button(tools, text="Deselect All",
                   style="Secondary.TButton",
                   command=self._deselect_all_vars).pack(side="right", padx=4)
        ttk.Button(tools, text="Auto-detect",
                   style="Secondary.TButton",
                   command=self._redetect_vars).pack(side="right", padx=4)

        # Search bar
        search_frame = tk.Frame(panel, bg=COLOR_PANEL)
        search_frame.pack(fill="x", pady=(0, 8))
        tk.Label(search_frame, text="Filter:",
                 font=FONT_BODY,
                 bg=COLOR_PANEL, fg=COLOR_TEXT_MUTED).pack(side="left",
                                                            padx=(0, 6))
        self.search_var = tk.StringVar()
        self.search_var.trace_add("write", lambda *_: self._apply_var_filter())
        entry = tk.Entry(search_frame, textvariable=self.search_var,
                         font=FONT_BODY, relief="solid", bd=1,
                         bg="#F9FAFC")
        entry.pack(side="left", fill="x", expand=True, ipady=4)

        # Treeview of variables
        tree_frame = tk.Frame(panel, bg=COLOR_PANEL)
        tree_frame.pack(fill="both", expand=True)

        self.vars_tree = ttk.Treeview(
            tree_frame,
            columns=("selected", "variable", "type", "completeness"),
            show="headings",
            selectmode="none")
        self.vars_tree.heading("selected", text="Use")
        self.vars_tree.heading("variable", text="Variable")
        self.vars_tree.heading("type", text="Type")
        self.vars_tree.heading("completeness", text="Complete")
        self.vars_tree.column("selected", width=60, anchor="center")
        self.vars_tree.column("variable", width=700, anchor="w")
        self.vars_tree.column("type", width=100, anchor="center")
        self.vars_tree.column("completeness", width=90, anchor="center")
        self.vars_tree.pack(side="left", fill="both", expand=True)

        scroll = ttk.Scrollbar(tree_frame, orient="vertical",
                               command=self.vars_tree.yview)
        scroll.pack(side="right", fill="y")
        self.vars_tree.configure(yscrollcommand=scroll.set)

        self.vars_tree.bind("<Button-1>", self._on_var_click)

        # Nav
        nav = tk.Frame(panel, bg=COLOR_PANEL)
        nav.pack(fill="x", pady=(14, 0))
        ttk.Button(nav, text="← Back",
                   style="Secondary.TButton",
                   command=lambda: self.notebook.select(self.tab_load)
                   ).pack(side="left")
        self.btn_to_run = ttk.Button(
            nav, text="Next: Run Analysis →",
            style="Primary.TButton",
            command=lambda: self.notebook.select(self.tab_run))
        self.btn_to_run.pack(side="right")

    # ---------------------------------------------------------------
    def _build_tab_run(self):
        tab = self.tab_run
        panel = ttk.Frame(tab, style="Panel.TFrame", padding=30)
        panel.pack(fill="both", expand=True, padx=16, pady=16)

        ttk.Label(panel, text="Step 3 — Run the analysis",
                  style="H2.TLabel").pack(anchor="w")
        ttk.Label(panel,
                  text="Configure clustering settings and generate segments.",
                  style="Muted.TLabel").pack(anchor="w", pady=(4, 20))

        # Options section
        opts_frame = tk.Frame(panel, bg=COLOR_PANEL)
        opts_frame.pack(fill="x", pady=(0, 20))

        # Number of segments
        row1 = tk.Frame(opts_frame, bg=COLOR_PANEL)
        row1.pack(fill="x", pady=6)
        tk.Label(row1, text="Number of segments (k):",
                 font=("Segoe UI", 10, "bold"),
                 bg=COLOR_PANEL, fg=COLOR_TEXT,
                 width=28, anchor="w").pack(side="left")

        self.k_mode_var = tk.StringVar(value="auto")
        tk.Radiobutton(row1, text="Auto-select (recommended)",
                       variable=self.k_mode_var, value="auto",
                       bg=COLOR_PANEL, fg=COLOR_TEXT,
                       selectcolor=COLOR_PANEL,
                       font=FONT_BODY,
                       command=self._update_k_mode).pack(side="left", padx=8)
        tk.Radiobutton(row1, text="Fixed:",
                       variable=self.k_mode_var, value="fixed",
                       bg=COLOR_PANEL, fg=COLOR_TEXT,
                       selectcolor=COLOR_PANEL,
                       font=FONT_BODY,
                       command=self._update_k_mode).pack(side="left", padx=8)

        self.k_spinbox = tk.Spinbox(row1, from_=2, to=10, width=5,
                                    font=FONT_BODY, state="disabled")
        self.k_spinbox.pack(side="left", padx=4)
        self.k_spinbox.delete(0, "end")
        self.k_spinbox.insert(0, "4")

        # Summary box
        self.run_summary = tk.Frame(panel, bg="#F0F4FA",
                                    relief="solid", bd=1)
        self.run_summary.pack(fill="x", pady=(10, 20))

        self.summary_inner = tk.Frame(self.run_summary, bg="#F0F4FA")
        self.summary_inner.pack(padx=20, pady=16, fill="x")

        # Progress
        self.progress_frame = tk.Frame(panel, bg=COLOR_PANEL)
        self.progress_frame.pack(fill="x", pady=10)

        self.progress_label = tk.Label(
            self.progress_frame, text="",
            font=FONT_BODY, bg=COLOR_PANEL, fg=COLOR_TEXT_MUTED)
        self.progress_label.pack(anchor="w")

        self.progress = ttk.Progressbar(
            self.progress_frame,
            style="Horizontal.TProgressbar",
            mode="determinate", length=400)
        self.progress.pack(fill="x", pady=(4, 0))

        # Nav
        nav = tk.Frame(panel, bg=COLOR_PANEL)
        nav.pack(fill="x", pady=(14, 0))
        ttk.Button(nav, text="← Back",
                   style="Secondary.TButton",
                   command=lambda: self.notebook.select(self.tab_vars)
                   ).pack(side="left")

        self.btn_run = ttk.Button(
            nav, text="▶  Run Segmentation",
            style="Primary.TButton",
            command=self._run_analysis)
        self.btn_run.pack(side="right")

    # ---------------------------------------------------------------
    def _build_tab_results(self):
        tab = self.tab_results
        panel = ttk.Frame(tab, style="Panel.TFrame", padding=30)
        panel.pack(fill="both", expand=True, padx=16, pady=16)

        ttk.Label(panel, text="Step 4 — Results & export",
                  style="H2.TLabel").pack(anchor="w")
        ttk.Label(panel,
                  text="Review the segments and export deliverables.",
                  style="Muted.TLabel").pack(anchor="w", pady=(4, 20))

        # Segments overview (cards)
        self.results_cards = tk.Frame(panel, bg=COLOR_PANEL)
        self.results_cards.pack(fill="both", expand=True)

        self.no_results_label = tk.Label(
            self.results_cards,
            text="Results will appear here after running the analysis.",
            font=FONT_BODY,
            bg=COLOR_PANEL, fg=COLOR_TEXT_MUTED)
        self.no_results_label.pack(pady=60)

        # Export buttons
        export_frame = tk.Frame(panel, bg=COLOR_PANEL)
        export_frame.pack(fill="x", pady=(20, 0))

        ttk.Label(export_frame, text="Export to:",
                  style="Panel.TLabel",
                  font=("Segoe UI", 10, "bold")).pack(side="left",
                                                       padx=(0, 12))

        self.btn_export_xlsx = ttk.Button(
            export_frame, text="📊  Excel Workbook",
            style="Primary.TButton",
            command=self._export_excel,
            state="disabled")
        self.btn_export_xlsx.pack(side="left", padx=4)

        self.btn_export_pptx = ttk.Button(
            export_frame, text="🎯  PowerPoint Deck",
            style="Primary.TButton",
            command=self._export_pptx,
            state="disabled")
        self.btn_export_pptx.pack(side="left", padx=4)

        self.btn_export_both = ttk.Button(
            export_frame, text="📦  Export Both",
            style="Primary.TButton",
            command=self._export_both,
            state="disabled")
        self.btn_export_both.pack(side="left", padx=4)

    # ---------------------------------------------------------------
    def _build_footer(self):
        footer = tk.Frame(self, bg=COLOR_PANEL, height=130,
                          highlightbackground=COLOR_BORDER,
                          highlightthickness=1)
        footer.pack(fill="x", side="bottom", padx=16, pady=(0, 14))
        footer.pack_propagate(False)

        tk.Label(footer, text="Activity Log",
                 font=("Segoe UI", 9, "bold"),
                 bg=COLOR_PANEL, fg=COLOR_PRIMARY).pack(
                     anchor="w", padx=10, pady=(6, 2))

        self.log_box = scrolledtext.ScrolledText(
            footer, height=5, font=FONT_MONO,
            bg="#F9FAFC", fg=COLOR_TEXT,
            relief="flat", bd=0, wrap="word")
        self.log_box.pack(fill="both", expand=True, padx=10, pady=(0, 8))

        self.log_box.tag_configure("info", foreground=COLOR_TEXT)
        self.log_box.tag_configure("success", foreground=COLOR_SUCCESS,
                                   font=("Consolas", 9, "bold"))
        self.log_box.tag_configure("warn", foreground=COLOR_WARN)
        self.log_box.tag_configure("error", foreground=COLOR_ERROR,
                                   font=("Consolas", 9, "bold"))
        self.log_box.configure(state="disabled")

        self._log("SegmentPro ready. Load a file to begin.", "info")

    # ---------------------------------------------------------------
    # ACTIONS
    # ---------------------------------------------------------------
    def _log(self, message, level="info"):
        """Thread-safe logger."""
        def do():
            ts = datetime.now().strftime("%H:%M:%S")
            self.log_box.configure(state="normal")
            self.log_box.insert("end", f"[{ts}] {message}\n", level)
            self.log_box.see("end")
            self.log_box.configure(state="disabled")
        try:
            self.after(0, do)
        except RuntimeError:
            pass

    def _choose_file(self):
        filetypes = [
            ("Data files", "*.xlsx *.xls *.xlsm *.csv"),
            ("Excel files", "*.xlsx *.xls *.xlsm"),
            ("CSV files", "*.csv"),
            ("All files", "*.*"),
        ]
        path = filedialog.askopenfilename(
            title="Select survey data file",
            filetypes=filetypes)
        if not path:
            return
        self.input_filepath = path
        self.file_label.configure(
            text=os.path.basename(path),
            fg=COLOR_PRIMARY)
        # Run load in a thread so UI stays responsive
        threading.Thread(target=self._do_load, daemon=True).start()

    def _do_load(self):
        try:
            self.engine.load_file(self.input_filepath)
            self.after(0, self._after_load_success)
        except Exception as ex:
            self._log(f"Error loading file: {ex}", "error")
            self.after(0, lambda: messagebox.showerror(
                "Load Error",
                f"Could not load file:\n\n{ex}"))

    def _after_load_success(self):
        df = self.engine.df
        # Build summary
        for w in self.dataset_summary.winfo_children():
            w.destroy()
        for w in self.preview_frame.winfo_children():
            w.destroy()

        # Summary stats
        stats_row = tk.Frame(self.dataset_summary, bg=COLOR_PANEL)
        stats_row.pack(fill="x")

        for label, val, color in [
            ("Rows", f"{len(df):,}", COLOR_PRIMARY),
            ("Columns", f"{len(df.columns):,}", COLOR_PRIMARY),
            ("Missing cells", f"{df.isna().sum().sum():,}", COLOR_WARN),
            ("File", os.path.basename(self.input_filepath), COLOR_TEXT_MUTED),
        ]:
            card = tk.Frame(stats_row, bg="#F0F4FA", padx=14, pady=10)
            card.pack(side="left", fill="y", padx=(0, 10))
            tk.Label(card, text=label, font=FONT_SMALL,
                     bg="#F0F4FA", fg=COLOR_TEXT_MUTED).pack(anchor="w")
            tk.Label(card, text=val,
                     font=("Segoe UI", 14, "bold"),
                     bg="#F0F4FA", fg=color).pack(anchor="w")

        # Preview table
        tk.Label(self.preview_frame, text="Preview (first 5 rows)",
                 font=FONT_SMALL, bg=COLOR_PANEL,
                 fg=COLOR_TEXT_MUTED).pack(anchor="w", pady=(14, 4))

        prev_frame = tk.Frame(self.preview_frame, bg=COLOR_PANEL)
        prev_frame.pack(fill="both", expand=True)

        cols_to_show = list(df.columns[:8])  # cap width
        tree = ttk.Treeview(prev_frame, columns=cols_to_show,
                            show="headings", height=5)
        for c in cols_to_show:
            tree.heading(c, text=str(c)[:40])
            tree.column(c, width=130, anchor="w")
        for _, row in df.head(5).iterrows():
            tree.insert("", "end",
                        values=[str(row[c])[:40] for c in cols_to_show])
        tree.pack(side="left", fill="both", expand=True)

        xscroll = ttk.Scrollbar(prev_frame, orient="horizontal",
                                command=tree.xview)
        xscroll.pack(side="bottom", fill="x")
        tree.configure(xscrollcommand=xscroll.set)

        if len(df.columns) > 8:
            tk.Label(self.preview_frame,
                     text=f"(showing first 8 of {len(df.columns)} columns)",
                     font=FONT_SMALL,
                     bg=COLOR_PANEL, fg=COLOR_TEXT_MUTED).pack(anchor="w",
                                                                pady=(4, 0))

        # Auto-detect variables and populate tab 2
        self.engine.detect_demographics()
        detected = self.engine.auto_detect_seg_cols()
        self._populate_vars(detected)
        self.btn_to_vars.configure(state="normal")
        self._update_tab_state()

    def _populate_vars(self, detected):
        self.vars_tree.delete(*self.vars_tree.get_children())
        self.selected_cols = set(detected)
        df = self.engine.df
        n = len(df)

        all_cols = list(df.columns)
        for col in all_cols:
            is_selected = col in self.selected_cols
            s = pd.to_numeric(df[col], errors="coerce")
            if s.notna().sum() / n > 0.5:
                unique = s.dropna().nunique()
                dtype = f"num ({unique}u)"
            else:
                dtype = "text"
            compl = f"{df[col].notna().sum() / n * 100:.0f}%"
            self.vars_tree.insert(
                "", "end",
                iid=col,
                values=("✓" if is_selected else "", str(col)[:200],
                        dtype, compl))
        self._update_var_count()
        self._log(
            f"Auto-detected {len(detected)} segmentation variables "
            f"(out of {len(all_cols)} total columns)",
            "success")

    def _on_var_click(self, event):
        region = self.vars_tree.identify("region", event.x, event.y)
        if region != "cell":
            return
        col = self.vars_tree.identify_column(event.x)
        row_id = self.vars_tree.identify_row(event.y)
        if not row_id:
            return
        # Toggle when user clicks the "Use" column
        if col == "#1":
            if row_id in self.selected_cols:
                self.selected_cols.discard(row_id)
                self.vars_tree.set(row_id, "selected", "")
            else:
                self.selected_cols.add(row_id)
                self.vars_tree.set(row_id, "selected", "✓")
            self._update_var_count()

    def _update_var_count(self):
        n = len(self.selected_cols)
        self.var_count_label.configure(text=f"{n} variables selected")
        if hasattr(self, "btn_to_run"):
            self.btn_to_run.configure(
                state="normal" if n >= 3 else "disabled")

    def _select_all_vars(self):
        for iid in self.vars_tree.get_children():
            self.selected_cols.add(iid)
            self.vars_tree.set(iid, "selected", "✓")
        self._update_var_count()

    def _deselect_all_vars(self):
        self.selected_cols.clear()
        for iid in self.vars_tree.get_children():
            self.vars_tree.set(iid, "selected", "")
        self._update_var_count()

    def _redetect_vars(self):
        detected = self.engine.auto_detect_seg_cols()
        self.selected_cols = set(detected)
        for iid in self.vars_tree.get_children():
            self.vars_tree.set(iid, "selected",
                               "✓" if iid in detected else "")
        self._update_var_count()
        self._log("Variables re-detected via auto-detection.", "info")

    def _apply_var_filter(self):
        query = self.search_var.get().lower().strip()
        # Rebuild tree items using saved full list
        df = self.engine.df
        if df is None:
            return
        self.vars_tree.delete(*self.vars_tree.get_children())
        n = len(df)
        for col in df.columns:
            if query and query not in str(col).lower():
                continue
            is_selected = col in self.selected_cols
            s = pd.to_numeric(df[col], errors="coerce")
            if s.notna().sum() / n > 0.5:
                unique = s.dropna().nunique()
                dtype = f"num ({unique}u)"
            else:
                dtype = "text"
            compl = f"{df[col].notna().sum() / n * 100:.0f}%"
            self.vars_tree.insert("", "end", iid=col,
                                  values=("✓" if is_selected else "",
                                          str(col)[:200], dtype, compl))

    # ---------------------------------------------------------------
    def _update_k_mode(self):
        if self.k_mode_var.get() == "fixed":
            self.k_spinbox.configure(state="normal")
        else:
            self.k_spinbox.configure(state="disabled")

    def _update_run_summary(self):
        for w in self.summary_inner.winfo_children():
            w.destroy()
        if self.engine.df is None:
            return
        tk.Label(self.summary_inner,
                 text="Analysis Summary",
                 font=("Segoe UI", 10, "bold"),
                 bg="#F0F4FA", fg=COLOR_PRIMARY).pack(anchor="w")

        n_records = len(self.engine.df)
        n_vars = len(self.selected_cols)
        k_text = ("auto-detect" if self.k_mode_var.get() == "auto"
                  else self.k_spinbox.get())
        lines = [
            f"• {n_records:,} records will be clustered",
            f"• {n_vars} segmentation variables selected",
            f"• k = {k_text}",
            f"• Algorithm: K-Means (standardized z-scores, 50 inits, seed=42)",
        ]
        for line in lines:
            tk.Label(self.summary_inner, text=line,
                     font=FONT_BODY, bg="#F0F4FA",
                     fg=COLOR_TEXT).pack(anchor="w", pady=1)

    # ---------------------------------------------------------------
    def _run_analysis(self):
        if self.is_running:
            return
        if len(self.selected_cols) < 3:
            messagebox.showwarning(
                "Too Few Variables",
                "Please select at least 3 variables for segmentation.")
            return
        self.is_running = True
        self.btn_run.configure(state="disabled")
        self.progress.configure(value=0)
        self.progress_label.configure(
            text="Starting analysis...", fg=COLOR_PRIMARY)
        threading.Thread(target=self._do_analysis, daemon=True).start()

    def _do_analysis(self):
        try:
            self._set_progress(10, "Cleaning data...")
            self.engine.clean_data(list(self.selected_cols))

            self._set_progress(30, "Evaluating cluster counts...")
            if self.k_mode_var.get() == "auto":
                k = None
            else:
                try:
                    k = int(self.k_spinbox.get())
                    if k < 2 or k > 10:
                        raise ValueError
                except ValueError:
                    self._log("Invalid k; using auto-detect.", "warn")
                    k = None

            self._set_progress(55, "Running K-Means clustering...")
            self.engine.run_clustering(k=k)

            self._set_progress(85, "Building segment profiles...")
            self._set_progress(100, "Analysis complete ✓")
            self.after(0, self._after_analysis_success)
        except Exception as ex:
            tb = traceback.format_exc()
            self._log(f"Analysis failed: {ex}\n{tb}", "error")
            self.after(0, lambda: messagebox.showerror(
                "Analysis Error",
                f"The analysis failed:\n\n{ex}"))
        finally:
            self.after(0, lambda: self.btn_run.configure(state="normal"))
            self.is_running = False

    def _set_progress(self, value, msg):
        def do():
            self.progress.configure(value=value)
            self.progress_label.configure(text=msg, fg=COLOR_PRIMARY)
        self.after(0, do)

    def _after_analysis_success(self):
        self._log(
            f"✓ Segmentation complete — {self.engine.k} segments identified.",
            "success")
        self._render_results()
        # Enable export buttons
        for btn in (self.btn_export_xlsx, self.btn_export_pptx,
                    self.btn_export_both):
            btn.configure(state="normal")
        # Jump to results tab
        self.notebook.select(self.tab_results)

    def _render_results(self):
        # Clear
        for w in self.results_cards.winfo_children():
            w.destroy()

        e = self.engine
        if e.labels is None:
            self.no_results_label = tk.Label(
                self.results_cards,
                text="Results will appear here after running the analysis.",
                font=FONT_BODY,
                bg=COLOR_PANEL, fg=COLOR_TEXT_MUTED)
            self.no_results_label.pack(pady=60)
            return

        sizes = e.segment_sizes()
        total = sum(sizes.values())

        # Summary strip
        strip = tk.Frame(self.results_cards, bg=COLOR_PANEL)
        strip.pack(fill="x", pady=(0, 16))
        for label, value, color in [
            ("Segments found", str(e.k), COLOR_PRIMARY),
            ("Respondents", f"{len(e.labels):,}", COLOR_PRIMARY),
            ("Variables used", str(len(e.seg_cols)), COLOR_PRIMARY),
            ("Best silhouette",
             f"{max(e.silhouette_scores.values()):.3f}"
             if e.silhouette_scores else "—",
             COLOR_SUCCESS),
        ]:
            card = tk.Frame(strip, bg="#F0F4FA", padx=16, pady=10)
            card.pack(side="left", fill="y", padx=(0, 10))
            tk.Label(card, text=label, font=FONT_SMALL,
                     bg="#F0F4FA", fg=COLOR_TEXT_MUTED).pack(anchor="w")
            tk.Label(card, text=value,
                     font=("Segoe UI", 16, "bold"),
                     bg="#F0F4FA", fg=color).pack(anchor="w")

        # Scroll area for cards
        scroll_frame = tk.Frame(self.results_cards, bg=COLOR_PANEL)
        scroll_frame.pack(fill="both", expand=True)

        canvas = tk.Canvas(scroll_frame, bg=COLOR_PANEL,
                           highlightthickness=0)
        canvas.pack(side="left", fill="both", expand=True)
        vsb = ttk.Scrollbar(scroll_frame, orient="vertical",
                            command=canvas.yview)
        vsb.pack(side="right", fill="y")
        canvas.configure(yscrollcommand=vsb.set)

        inner = tk.Frame(canvas, bg=COLOR_PANEL)
        inner_window_id = canvas.create_window((0, 0), window=inner, anchor="nw")

        def on_frame_configure(_):
            canvas.configure(scrollregion=canvas.bbox("all"))
        inner.bind("<Configure>", on_frame_configure)

        def on_canvas_configure(event):
            # Stretch the inner frame to match the canvas width so cards fill
            # the available horizontal space.
            canvas.itemconfigure(inner_window_id, width=event.width)
        canvas.bind("<Configure>", on_canvas_configure)

        def on_mousewheel(event):
            canvas.yview_scroll(int(-1 * (event.delta / 120)), "units")
        canvas.bind_all("<MouseWheel>", on_mousewheel)

        for seg in sorted(sizes.keys()):
            self._build_segment_card(inner, seg, sizes[seg], total)

    def _build_segment_card(self, parent, seg, size, total):
        e = self.engine
        card = tk.Frame(parent, bg="#FFFFFF",
                        highlightbackground=COLOR_BORDER,
                        highlightthickness=1)
        card.pack(fill="x", pady=6)

        # Top strip
        strip = tk.Frame(card, bg=COLOR_PRIMARY, height=6)
        strip.pack(fill="x")

        content = tk.Frame(card, bg="#FFFFFF")
        content.pack(fill="x", padx=20, pady=14)

        # Header row
        header = tk.Frame(content, bg="#FFFFFF")
        header.pack(fill="x")

        # Badge
        badge = tk.Frame(header, bg=COLOR_PRIMARY, width=50, height=50)
        badge.pack(side="left")
        badge.pack_propagate(False)
        tk.Label(badge, text=str(seg),
                 font=("Georgia", 22, "bold"),
                 bg=COLOR_PRIMARY, fg="#FFFFFF").pack(
                     expand=True, fill="both")

        info = tk.Frame(header, bg="#FFFFFF")
        info.pack(side="left", padx=14, fill="x", expand=True)

        tk.Label(info,
                 text=e.seg_names.get(seg, f"Segment {seg}"),
                 font=("Segoe UI", 14, "bold"),
                 bg="#FFFFFF", fg=COLOR_PRIMARY).pack(anchor="w")
        pct = size / total * 100
        tk.Label(info,
                 text=f"n = {size:,}  ·  {pct:.1f}% of respondents",
                 font=FONT_SMALL,
                 bg="#FFFFFF", fg=COLOR_TEXT_MUTED).pack(anchor="w")

        # Description
        desc = e.seg_descriptions.get(seg, "")
        tk.Label(content, text=desc, wraplength=950, justify="left",
                 font=FONT_BODY, bg="#FFFFFF", fg=COLOR_TEXT).pack(
                     anchor="w", pady=(10, 4))

        # Top 3 differentiators
        z = e.seg_std_profile.loc[seg]
        top3 = z.abs().sort_values(ascending=False).head(3)
        tk.Label(content, text="Key differentiators:",
                 font=("Segoe UI", 9, "bold"),
                 bg="#FFFFFF", fg=COLOR_TEXT_MUTED).pack(anchor="w",
                                                          pady=(8, 2))
        for var, _ in top3.items():
            zval = z[var]
            direction = "↑" if zval > 0 else "↓"
            color = COLOR_SUCCESS if zval > 0 else COLOR_ERROR
            row = tk.Frame(content, bg="#FFFFFF")
            row.pack(fill="x", pady=1)
            # Pack fixed-width items first so tkinter reserves their space
            # before the expanding name label.
            tk.Label(row, text=direction, font=("Segoe UI", 11, "bold"),
                     bg="#FFFFFF", fg=color, width=2).pack(side="left")
            tk.Label(row, text=f"{zval:+.2f}σ",
                     font=("Consolas", 9, "bold"),
                     bg="#FFFFFF", fg=color, width=8,
                     anchor="e").pack(side="right", padx=(6, 0))
            tk.Label(row, text=smart_short(var, max_len=95),
                     font=FONT_SMALL, bg="#FFFFFF", fg=COLOR_TEXT,
                     anchor="w", justify="left").pack(side="left",
                                                       fill="x", expand=True)

    # ---------------------------------------------------------------
    def _choose_output_dir(self):
        path = filedialog.askdirectory(title="Choose output folder")
        if path:
            self.output_dir = path
        return path

    def _default_output_dir(self):
        if self.output_dir:
            return self.output_dir
        if self.input_filepath:
            return os.path.dirname(self.input_filepath)
        return os.getcwd()

    def _timestamped_name(self, base, ext):
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        return f"{base}_{stamp}.{ext}"

    def _export_excel(self):
        path = filedialog.asksaveasfilename(
            title="Save Excel report",
            defaultextension=".xlsx",
            initialdir=self._default_output_dir(),
            initialfile=self._timestamped_name("Segments_Report", "xlsx"),
            filetypes=[("Excel Workbook", "*.xlsx")])
        if not path:
            return
        threading.Thread(target=self._do_export_excel, args=(path,),
                         daemon=True).start()

    def _do_export_excel(self, path):
        try:
            gen = ReportGenerator(self.engine, log_callback=self._log)
            gen.write_excel(path)
            self.after(0, lambda: messagebox.showinfo(
                "Export Complete",
                f"Excel report saved to:\n\n{path}"))
        except Exception as ex:
            self._log(f"Excel export failed: {ex}", "error")
            self.after(0, lambda: messagebox.showerror(
                "Export Error",
                f"Could not save Excel file:\n\n{ex}"))

    def _export_pptx(self):
        path = filedialog.asksaveasfilename(
            title="Save PowerPoint deck",
            defaultextension=".pptx",
            initialdir=self._default_output_dir(),
            initialfile=self._timestamped_name("Segments_Deck", "pptx"),
            filetypes=[("PowerPoint Presentation", "*.pptx")])
        if not path:
            return
        threading.Thread(target=self._do_export_pptx, args=(path,),
                         daemon=True).start()

    def _do_export_pptx(self, path):
        try:
            gen = ReportGenerator(self.engine, log_callback=self._log)
            gen.write_pptx(path)
            self.after(0, lambda: messagebox.showinfo(
                "Export Complete",
                f"PowerPoint deck saved to:\n\n{path}"))
        except ImportError:
            self.after(0, lambda: messagebox.showerror(
                "Missing Package",
                "PowerPoint export requires the 'python-pptx' package.\n\n"
                "Install it with:\n   pip install python-pptx"))
        except Exception as ex:
            self._log(f"PowerPoint export failed: {ex}", "error")
            self.after(0, lambda: messagebox.showerror(
                "Export Error",
                f"Could not save PowerPoint file:\n\n{ex}"))

    def _export_both(self):
        outdir = self._choose_output_dir()
        if not outdir:
            return
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        xlsx_path = os.path.join(outdir, f"Segments_Report_{stamp}.xlsx")
        pptx_path = os.path.join(outdir, f"Segments_Deck_{stamp}.pptx")
        threading.Thread(target=self._do_export_both,
                         args=(xlsx_path, pptx_path),
                         daemon=True).start()

    def _do_export_both(self, xlsx_path, pptx_path):
        try:
            gen = ReportGenerator(self.engine, log_callback=self._log)
            gen.write_excel(xlsx_path)
            gen.write_pptx(pptx_path)
            msg = (f"Both files saved to:\n\n"
                   f"• {os.path.basename(xlsx_path)}\n"
                   f"• {os.path.basename(pptx_path)}\n\n"
                   f"Folder:\n{os.path.dirname(xlsx_path)}")
            self.after(0, lambda: messagebox.showinfo(
                "Export Complete", msg))
        except Exception as ex:
            self._log(f"Export failed: {ex}", "error")
            self.after(0, lambda: messagebox.showerror(
                "Export Error",
                f"Could not save all files:\n\n{ex}"))

    # ---------------------------------------------------------------
    def _update_tab_state(self):
        # When switching to the run tab, refresh its summary
        def on_tab_change(_):
            if self.notebook.index("current") == 2:
                self._update_run_summary()
        self.notebook.bind("<<NotebookTabChanged>>", on_tab_change)


# ============================================================================
# MAIN
# ============================================================================
def main():
    if not _TK_AVAILABLE:
        print("ERROR: Tkinter is not available in this Python installation.\n"
              "On Windows, re-install Python with the 'tcl/tk and IDLE' option.")
        sys.exit(1)

    # Dependency check
    missing = []
    try:
        import sklearn
    except ImportError:
        missing.append("scikit-learn")
    try:
        import openpyxl
    except ImportError:
        missing.append("openpyxl")
    try:
        import pptx
    except ImportError:
        missing.append("python-pptx")

    if missing:
        root = tk.Tk()
        root.withdraw()
        messagebox.showerror(
            f"{APP_NAME} — Missing Dependencies",
            f"The following Python packages are required:\n\n"
            f"  {', '.join(missing)}\n\n"
            f"Install them with:\n"
            f"  pip install {' '.join(missing)}")
        return

    app = SegmentProApp()
    app.mainloop()


if __name__ == "__main__":
    main()
